import openpyxl
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from reportlab.lib.pagesizes import letter, A4
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.lib.colors import HexColor, black, white
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak, Image
from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_JUSTIFY
import requests
import os
from typing import Dict, List, Optional
from datetime import datetime
import json
import asyncio
from groq import Groq


class DeliverableService:
    """Service for generating professional consulting deliverables (Excel, PPT, PDF)"""

    # Professional color scheme
    COLORS = {
        'primary': '0F172A',      # Deep Navy
        'secondary': '1E293B',    # Slate
        'accent': '2563EB',       # Electric Blue
        'accent_light': '3B82F6', # Light Blue
        'text': 'F8FAFC',         # White text
        'text_dark': '1E293B',    # Dark text
        'success': '10B981',      # Green
        'warning': 'F59E0B',      # Amber
    }

    def __init__(self):
        self.gamma_api_key = os.getenv('GAMMA_API_KEY', '')
        self.groq_api_key = os.getenv('GROQ_API_KEYS', '').split(',')[0].strip()
        self.deliverables_path = "/home/user/aiconsultant/deliverables"
        os.makedirs(self.deliverables_path, exist_ok=True)

    async def verify_content_online(self, claims: List[str], context: str) -> Dict:
        """Verify and enrich content using AI with web knowledge"""
        if not self.groq_api_key:
            return {"verified": False, "enriched_claims": claims}

        try:
            client = Groq(api_key=self.groq_api_key)

            claims_text = "\n".join([f"- {c}" for c in claims[:10]])

            response = client.chat.completions.create(
                messages=[{
                    "role": "user",
                    "content": f"""As a consulting analyst, verify and enrich these claims with current market data:

Context: {context}

Claims to verify:
{claims_text}

For each claim, provide:
1. Verification status (Verified/Partially Verified/Needs Update)
2. Current accurate data or statistics (with approximate year)
3. Any important updates or corrections

Format as JSON with "verified_claims" array containing objects with "original", "status", "enriched", "source_note" fields."""
                }],
                model="llama-3.3-70b-versatile",
                max_tokens=1500,
                temperature=0.3
            )

            result_text = response.choices[0].message.content

            # Try to parse JSON from response
            try:
                if '{' in result_text:
                    json_start = result_text.index('{')
                    json_end = result_text.rindex('}') + 1
                    result = json.loads(result_text[json_start:json_end])
                    return {"verified": True, "data": result}
            except:
                pass

            return {"verified": True, "enriched_text": result_text}

        except Exception as e:
            return {"verified": False, "error": str(e), "enriched_claims": claims}

    async def generate_ai_slide_content(self, topic: str, context: str, slide_type: str) -> Dict:
        """Generate rich slide content using AI"""
        if not self.groq_api_key:
            return None

        try:
            client = Groq(api_key=self.groq_api_key)

            prompts = {
                "executive_summary": f"Create a compelling executive summary for: {topic}. Context: {context}. Provide 3-4 key bullet points.",
                "market_analysis": f"Provide current market analysis for: {topic}. Include market size, growth rate, and key trends. Context: {context}",
                "competitive_landscape": f"Analyze the competitive landscape for: {topic}. List top 3-4 competitors with their strengths. Context: {context}",
                "recommendations": f"Provide strategic recommendations for: {topic}. List 4-5 actionable items with expected impact. Context: {context}",
                "implementation": f"Create an implementation roadmap for: {topic}. Include phases, timeline, and key milestones. Context: {context}",
                "risks": f"Identify key risks and mitigation strategies for: {topic}. List top 4 risks with mitigation approaches. Context: {context}",
                "financial": f"Provide financial projections/considerations for: {topic}. Include investment needs, ROI expectations, and key metrics. Context: {context}"
            }

            prompt = prompts.get(slide_type, f"Provide insights on {slide_type} for: {topic}. Context: {context}")

            response = client.chat.completions.create(
                messages=[{"role": "user", "content": prompt + "\n\nKeep response concise with clear bullet points."}],
                model="llama-3.1-8b-instant",
                max_tokens=500,
                temperature=0.5
            )

            return {"content": response.choices[0].message.content, "type": slide_type}

        except Exception as e:
            return None

    # ==================== EXCEL GENERATION ====================

    def generate_excel_report(self, project_data: Dict, analysis_type: str) -> str:
        """Generate Excel financial model or analysis report"""
        wb = openpyxl.Workbook()

        if analysis_type == "financial_model":
            self._create_financial_model(wb, project_data)
        elif analysis_type == "current_state":
            self._create_current_state_analysis(wb, project_data)
        elif analysis_type == "future_state":
            self._create_future_state_analysis(wb, project_data)
        else:
            self._create_summary_report(wb, project_data)

        filename = f"{analysis_type}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.xlsx"
        filepath = os.path.join(self.deliverables_path, filename)
        wb.save(filepath)

        return filename

    def _create_financial_model(self, wb, data: Dict):
        """Create financial model worksheet"""
        ws = wb.active
        ws.title = "Financial Model"

        header_fill = PatternFill(start_color=self.COLORS['primary'], end_color=self.COLORS['primary'], fill_type="solid")
        header_font = Font(color=self.COLORS['text'], bold=True, size=12)
        accent_fill = PatternFill(start_color=self.COLORS['accent'], end_color=self.COLORS['accent'], fill_type="solid")

        ws['A1'] = data.get('title', 'Financial Analysis Model')
        ws['A1'].font = Font(bold=True, size=18, color=self.COLORS['primary'])
        ws.merge_cells('A1:E1')

        ws['A3'] = "Category"
        ws['B3'] = "Year 1"
        ws['C3'] = "Year 2"
        ws['D3'] = "Year 3"
        ws['E3'] = "Total"

        for cell in ['A3', 'B3', 'C3', 'D3', 'E3']:
            ws[cell].fill = header_fill
            ws[cell].font = header_font
            ws[cell].alignment = Alignment(horizontal='center')

        categories = ["Revenue", "Cost of Goods Sold", "Gross Profit", "Operating Expenses", "EBITDA", "Net Income"]

        for idx, category in enumerate(categories, start=4):
            ws[f'A{idx}'] = category
            ws[f'A{idx}'].font = Font(bold=True)

        ws.column_dimensions['A'].width = 25
        for col in ['B', 'C', 'D', 'E']:
            ws.column_dimensions[col].width = 15

    def _create_current_state_analysis(self, wb, data: Dict):
        """Create current state analysis worksheet"""
        ws = wb.active
        ws.title = "Current State"

        header_fill = PatternFill(start_color=self.COLORS['primary'], end_color=self.COLORS['primary'], fill_type="solid")

        ws['A1'] = data.get('title', 'Current State Analysis')
        ws['A1'].font = Font(bold=True, size=18, color=self.COLORS['primary'])

        ws['A3'] = "Problem Statement:"
        ws['A3'].font = Font(bold=True, size=12, color=self.COLORS['accent'])
        ws['A4'] = data.get('problem', 'No problem statement provided')
        ws['A4'].alignment = Alignment(wrap_text=True)

        ws['A6'] = "Key Findings:"
        ws['A6'].font = Font(bold=True, size=12, color=self.COLORS['accent'])

        findings = data.get('findings', [])
        for idx, finding in enumerate(findings, start=7):
            ws[f'A{idx}'] = f"• {finding}"
            ws[f'A{idx}'].alignment = Alignment(wrap_text=True)

        ws.column_dimensions['A'].width = 100

    def _create_future_state_analysis(self, wb, data: Dict):
        """Create future state analysis worksheet"""
        ws = wb.active
        ws.title = "Future State"

        ws['A1'] = data.get('title', 'Future State Vision')
        ws['A1'].font = Font(bold=True, size=18, color=self.COLORS['primary'])

        ws['A3'] = "Recommended Solution:"
        ws['A3'].font = Font(bold=True, size=12, color=self.COLORS['accent'])
        ws['A4'] = data.get('solution', 'No solution provided')
        ws['A4'].alignment = Alignment(wrap_text=True)

        ws['A6'] = "Implementation Roadmap:"
        ws['A6'].font = Font(bold=True, size=12, color=self.COLORS['accent'])

        roadmap = data.get('roadmap', [])
        for idx, step in enumerate(roadmap, start=7):
            ws[f'A{idx}'] = f"{idx-6}. {step}"
            ws[f'A{idx}'].alignment = Alignment(wrap_text=True)

        ws.column_dimensions['A'].width = 100

    def _create_summary_report(self, wb, data: Dict):
        """Create summary report worksheet"""
        ws = wb.active
        ws.title = "Executive Summary"

        ws['A1'] = "Executive Summary"
        ws['A1'].font = Font(bold=True, size=18, color=self.COLORS['primary'])

        ws['A3'] = "Project Overview:"
        ws['A3'].font = Font(bold=True, size=12, color=self.COLORS['accent'])
        ws['A4'] = data.get('overview', 'Summary not available')

        ws.column_dimensions['A'].width = 100

    # ==================== POWERPOINT GENERATION ====================

    async def generate_ppt(self, content: Dict) -> Dict:
        """Generate professional PowerPoint presentation"""
        try:
            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)

            title = content.get('title', 'Strategic Analysis Report')
            problem = content.get('problem', '')
            context = content.get('context', problem)
            findings = content.get('findings', [])
            recommendations = content.get('recommendations', [])
            debate_data = content.get('debate_data', {})
            research_data = content.get('research_data', {})

            # Generate AI-enriched content for each section
            sections_to_generate = ['executive_summary', 'market_analysis', 'recommendations', 'implementation', 'risks']
            ai_content = {}

            for section in sections_to_generate:
                result = await self.generate_ai_slide_content(title, context, section)
                if result:
                    ai_content[section] = result

            # Slide 1: Title Slide
            self._add_title_slide(prs, title, "Strategic Consulting Analysis")

            # Slide 2: Executive Summary
            exec_summary = ai_content.get('executive_summary', {}).get('content', '')
            if not exec_summary and findings:
                exec_summary = "\n".join([f"• {f}" for f in findings[:4]])
            self._add_content_slide(prs, "Executive Summary", self._parse_bullets(exec_summary) or findings[:5])

            # Slide 3: The Challenge
            self._add_content_slide(prs, "The Challenge", [problem] if problem else ["Business challenge analysis pending"])

            # Slide 4: Market Analysis (AI-generated)
            market_content = ai_content.get('market_analysis', {}).get('content', '')
            if market_content:
                self._add_content_slide(prs, "Market Analysis", self._parse_bullets(market_content))

            # Slide 5: Key Findings
            if findings:
                self._add_content_slide(prs, "Key Findings", findings[:6])

            # Slide 6: Consultant Perspectives (from debate)
            if debate_data:
                perspectives = []
                debate_history = debate_data.get('debate_history', [])
                seen_agents = set()
                for entry in debate_history:
                    agent = entry.get('agent', '')
                    if agent not in seen_agents and not entry.get('argument', '').startswith('['):
                        perspectives.append(f"{agent}: {entry.get('argument', '')[:100]}...")
                        seen_agents.add(agent)
                        if len(perspectives) >= 4:
                            break
                if perspectives:
                    self._add_content_slide(prs, "Expert Perspectives", perspectives)

            # Slide 7: Competitive Analysis (from research)
            if research_data:
                vendor_comparison = research_data.get('vendor_comparison', {})
                if vendor_comparison:
                    comp_items = []
                    for vendor, details in list(vendor_comparison.items())[:4]:
                        score = details.get('total_score', 0)
                        comp_items.append(f"{vendor}: {score}/10 - {details.get('best_for', 'General use')[:50]}")
                    self._add_content_slide(prs, "Competitive Landscape", comp_items)

            # Slide 8: Strategic Recommendations
            rec_content = ai_content.get('recommendations', {}).get('content', '')
            recs = self._parse_bullets(rec_content) if rec_content else recommendations[:5]
            if recs:
                self._add_content_slide(prs, "Strategic Recommendations", recs)

            # Slide 9: Implementation Roadmap
            impl_content = ai_content.get('implementation', {}).get('content', '')
            if impl_content:
                self._add_content_slide(prs, "Implementation Roadmap", self._parse_bullets(impl_content))

            # Slide 10: Risks & Mitigation
            risk_content = ai_content.get('risks', {}).get('content', '')
            if risk_content:
                self._add_content_slide(prs, "Risks & Mitigation", self._parse_bullets(risk_content))

            # Slide 11: Next Steps
            next_steps = content.get('next_steps', [
                "Review and validate recommendations with leadership team",
                "Develop detailed implementation plan with timelines",
                "Identify quick wins for immediate action",
                "Establish KPIs and success metrics",
                "Schedule follow-up review in 2 weeks"
            ])
            self._add_content_slide(prs, "Next Steps", next_steps[:5])

            # Slide 12: Thank You
            self._add_thank_you_slide(prs)

            filename = f"presentation_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
            filepath = os.path.join(self.deliverables_path, filename)
            prs.save(filepath)

            return {
                "success": True,
                "filename": filename,
                "filepath": filepath,
                "slides_count": len(prs.slides),
                "type": "pptx"
            }

        except Exception as e:
            return {"success": False, "error": str(e)}

    def _add_title_slide(self, prs, title: str, subtitle: str):
        """Add professional title slide"""
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)

        # Background shape
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(0x0F, 0x17, 0x2A)
        bg.line.fill.background()

        # Accent bar
        accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(3), prs.slide_width, Inches(0.1))
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = RGBColor(0x25, 0x63, 0xEB)
        accent_bar.line.fill.background()

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12), Inches(1.5))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xF8, 0xFA, 0xFC)
        p.alignment = PP_ALIGN.CENTER

        # Subtitle
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(12), Inches(0.8))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(0x94, 0xA3, 0xB8)
        p.alignment = PP_ALIGN.CENTER

        # Date
        date_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12), Inches(0.5))
        tf = date_box.text_frame
        p = tf.paragraphs[0]
        p.text = datetime.now().strftime("%B %Y")
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(0x64, 0x74, 0x8B)
        p.alignment = PP_ALIGN.CENTER

    def _add_content_slide(self, prs, title: str, bullet_points: List[str]):
        """Add content slide with bullet points"""
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)

        # Background
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(0xF8, 0xFA, 0xFC)
        bg.line.fill.background()

        # Header bar
        header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
        header.fill.solid()
        header.fill.fore_color.rgb = RGBColor(0x0F, 0x17, 0x2A)
        header.line.fill.background()

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xF8, 0xFA, 0xFC)

        # Content area
        content_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.5), Inches(5))
        tf = content_box.text_frame
        tf.word_wrap = True

        for i, point in enumerate(bullet_points[:7]):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            # Clean the point text
            point_text = str(point).strip()
            if point_text.startswith(('•', '-', '*', '→')):
                point_text = point_text[1:].strip()

            p.text = f"→  {point_text}"
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(0x1E, 0x29, 0x3B)
            p.space_before = Pt(12)
            p.space_after = Pt(8)

    def _add_thank_you_slide(self, prs):
        """Add thank you slide"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        # Background
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(0x0F, 0x17, 0x2A)
        bg.line.fill.background()

        # Thank you text
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(12), Inches(1.5))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "Thank You"
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0x25, 0x63, 0xEB)
        p.alignment = PP_ALIGN.CENTER

        # Subtitle
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(12), Inches(0.8))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = "Questions & Discussion"
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(0x94, 0xA3, 0xB8)
        p.alignment = PP_ALIGN.CENTER

    def _parse_bullets(self, text: str) -> List[str]:
        """Parse text into bullet points"""
        if not text:
            return []

        lines = text.split('\n')
        bullets = []
        for line in lines:
            line = line.strip()
            if line and len(line) > 5:
                # Remove common bullet prefixes
                for prefix in ['•', '-', '*', '→', '1.', '2.', '3.', '4.', '5.', '6.', '7.', '8.', '9.']:
                    if line.startswith(prefix):
                        line = line[len(prefix):].strip()
                        break
                if line:
                    bullets.append(line[:200])  # Limit length
        return bullets[:8]

    # ==================== PDF GENERATION ====================

    async def generate_pdf(self, content: Dict) -> Dict:
        """Generate professional PDF report"""
        try:
            filename = f"report_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pdf"
            filepath = os.path.join(self.deliverables_path, filename)

            doc = SimpleDocTemplate(
                filepath,
                pagesize=A4,
                rightMargin=50,
                leftMargin=50,
                topMargin=50,
                bottomMargin=50
            )

            styles = getSampleStyleSheet()

            # Custom styles
            styles.add(ParagraphStyle(
                name='CustomTitle',
                parent=styles['Title'],
                fontSize=28,
                textColor=HexColor('#0F172A'),
                spaceAfter=30,
                alignment=TA_CENTER
            ))

            styles.add(ParagraphStyle(
                name='CustomHeading',
                parent=styles['Heading1'],
                fontSize=18,
                textColor=HexColor('#2563EB'),
                spaceBefore=20,
                spaceAfter=12
            ))

            styles.add(ParagraphStyle(
                name='CustomBody',
                parent=styles['Normal'],
                fontSize=11,
                textColor=HexColor('#1E293B'),
                alignment=TA_JUSTIFY,
                spaceAfter=8,
                leading=16
            ))

            styles.add(ParagraphStyle(
                name='CustomBullet',
                parent=styles['Normal'],
                fontSize=11,
                textColor=HexColor('#1E293B'),
                leftIndent=20,
                spaceAfter=6,
                bulletIndent=10
            ))

            story = []

            # Title
            title = content.get('title', 'Strategic Analysis Report')
            story.append(Paragraph(title, styles['CustomTitle']))
            story.append(Spacer(1, 20))

            # Executive Summary
            story.append(Paragraph("Executive Summary", styles['CustomHeading']))
            problem = content.get('problem', 'Strategic analysis and recommendations')
            story.append(Paragraph(problem, styles['CustomBody']))
            story.append(Spacer(1, 15))

            # Key Findings
            findings = content.get('findings', [])
            if findings:
                story.append(Paragraph("Key Findings", styles['CustomHeading']))
                for finding in findings[:8]:
                    story.append(Paragraph(f"• {finding}", styles['CustomBullet']))
                story.append(Spacer(1, 15))

            # Debate Insights
            debate_data = content.get('debate_data', {})
            if debate_data:
                story.append(Paragraph("Expert Analysis", styles['CustomHeading']))

                consensus = debate_data.get('consensus', {})
                key_insights = consensus.get('key_insights', [])
                for insight in key_insights[:5]:
                    story.append(Paragraph(f"• {insight}", styles['CustomBullet']))
                story.append(Spacer(1, 15))

                # Agent perspectives
                story.append(Paragraph("Consultant Perspectives", styles['CustomHeading']))
                debate_history = debate_data.get('debate_history', [])
                seen_agents = set()
                for entry in debate_history[:16]:
                    agent = entry.get('agent', '')
                    argument = entry.get('argument', '')
                    if agent not in seen_agents and not argument.startswith('[') and len(argument) > 50:
                        story.append(Paragraph(f"<b>{agent}</b> ({entry.get('role', '')})", styles['CustomBody']))
                        story.append(Paragraph(argument[:400] + "..." if len(argument) > 400 else argument, styles['CustomBullet']))
                        story.append(Spacer(1, 10))
                        seen_agents.add(agent)
                        if len(seen_agents) >= 6:
                            break

            # Research Data
            research_data = content.get('research_data', {})
            if research_data:
                story.append(PageBreak())
                story.append(Paragraph("Market Research & Analysis", styles['CustomHeading']))

                analysis = research_data.get('analysis', {})
                market_position = analysis.get('market_position', '')
                if market_position:
                    story.append(Paragraph(market_position, styles['CustomBody']))

                # Vendor comparison table
                vendor_comparison = research_data.get('vendor_comparison', {})
                if vendor_comparison:
                    story.append(Spacer(1, 15))
                    story.append(Paragraph("Competitive Analysis", styles['CustomHeading']))

                    table_data = [['Vendor', 'Score', 'Best For', 'Pricing']]
                    for vendor, details in list(vendor_comparison.items())[:5]:
                        table_data.append([
                            vendor,
                            f"{details.get('total_score', 0)}/10",
                            details.get('best_for', 'N/A')[:40],
                            details.get('pricing_tier', 'N/A')
                        ])

                    table = Table(table_data, colWidths=[100, 60, 200, 80])
                    table.setStyle(TableStyle([
                        ('BACKGROUND', (0, 0), (-1, 0), HexColor('#0F172A')),
                        ('TEXTCOLOR', (0, 0), (-1, 0), white),
                        ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                        ('FONTSIZE', (0, 0), (-1, 0), 10),
                        ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                        ('BACKGROUND', (0, 1), (-1, -1), HexColor('#F8FAFC')),
                        ('FONTNAME', (0, 1), (-1, -1), 'Helvetica'),
                        ('FONTSIZE', (0, 1), (-1, -1), 9),
                        ('GRID', (0, 0), (-1, -1), 1, HexColor('#E2E8F0')),
                        ('VALIGN', (0, 0), (-1, -1), 'TOP'),
                    ]))
                    story.append(table)

            # Recommendations
            recommendations = content.get('recommendations', [])
            if recommendations:
                story.append(Spacer(1, 20))
                story.append(Paragraph("Strategic Recommendations", styles['CustomHeading']))
                for i, rec in enumerate(recommendations[:6], 1):
                    story.append(Paragraph(f"{i}. {rec}", styles['CustomBullet']))

            # Next Steps
            next_steps = content.get('next_steps', [])
            if next_steps:
                story.append(Spacer(1, 20))
                story.append(Paragraph("Next Steps", styles['CustomHeading']))
                for step in next_steps[:5]:
                    story.append(Paragraph(f"→ {step}", styles['CustomBullet']))

            # Footer info
            story.append(Spacer(1, 40))
            story.append(Paragraph(
                f"<i>Generated on {datetime.now().strftime('%B %d, %Y')} | AI Consultant Platform</i>",
                ParagraphStyle('Footer', parent=styles['Normal'], fontSize=9, textColor=HexColor('#64748B'), alignment=TA_CENTER)
            ))

            doc.build(story)

            return {
                "success": True,
                "filename": filename,
                "filepath": filepath,
                "type": "pdf"
            }

        except Exception as e:
            return {"success": False, "error": str(e)}

    # ==================== LEGACY METHODS (for backward compatibility) ====================

    async def generate_ppt_via_gamma(self, content: Dict) -> Dict:
        """Generate professional presentation using Gamma API"""
        if not self.gamma_api_key:
            # Fallback to python-pptx if no Gamma API key
            return await self.generate_ppt(content)

        try:
            title = content.get('title', 'Consulting Report')
            text = content.get('text', '')
            debate_data = content.get('debate_data', {})

            # Build rich content for Gamma
            presentation_content = f"""# {title}

## Executive Summary
{debate_data.get('consensus', {}).get('summary', 'Strategic consulting analysis and recommendations.')}

## The Challenge
{text[:1500] if text else 'Business challenge requiring strategic analysis.'}

## Key Insights
{chr(10).join(['• ' + insight for insight in debate_data.get('consensus', {}).get('key_insights', ['Market analysis completed', 'Strategic opportunities identified', 'Risk assessment performed'])[:5]])}

## Recommendations
{chr(10).join(['• ' + rec for rec in debate_data.get('consensus', {}).get('recommendations', ['Implement phased approach', 'Focus on core competencies', 'Monitor KPIs regularly'])[:5]])}

## Implementation Roadmap
{chr(10).join(debate_data.get('consensus', {}).get('next_steps', ['Phase 1: Planning and Assessment', 'Phase 2: Implementation', 'Phase 3: Monitoring and Optimization'])[:5])}

## Risk Factors
{chr(10).join(['• ' + risk for risk in debate_data.get('consensus', {}).get('risk_factors', ['Market volatility', 'Resource constraints', 'Timeline pressures'])[:4]])}

## Expected Impact
• Growth Potential: {debate_data.get('consensus', {}).get('estimated_impact', {}).get('arr_growth', '15-25% improvement')}
• Efficiency Gains: {debate_data.get('consensus', {}).get('estimated_impact', {}).get('cac_reduction', '20-30% optimization')}
• Timeline: {debate_data.get('consensus', {}).get('estimated_impact', {}).get('timeline', '6-12 months')}

## Next Steps
• Schedule follow-up meeting
• Begin Phase 1 implementation
• Establish monitoring framework
"""

            # Call Gamma API
            headers = {
                "Content-Type": "application/json",
                "X-API-KEY": self.gamma_api_key
            }

            payload = {
                "inputText": presentation_content,
                "textMode": "generate",
                "format": "presentation",
                "numCards": 10,
                "cardSplit": "auto",
                "cardOptions": {
                    "dimensions": "16x9"
                },
                "textOptions": {
                    "amount": "medium",
                    "tone": "professional",
                    "audience": "executives"
                },
                "imageOptions": {
                    "source": "unsplash",
                    "style": "professional"
                },
                "additionalInstructions": "Create a professional consulting presentation with clean design, executive-friendly visuals, and clear data visualization. Use a modern corporate color scheme."
            }

            # Start generation
            response = requests.post(
                "https://public-api.gamma.app/v1.0/generations",
                headers=headers,
                json=payload,
                timeout=30
            )

            if response.status_code != 200 and response.status_code != 201:
                print(f"Gamma API error: {response.status_code} - {response.text}")
                return await self.generate_ppt(content)  # Fallback

            result = response.json()
            generation_id = result.get('generationId') or result.get('id')

            if not generation_id:
                return await self.generate_ppt(content)  # Fallback

            # Poll for completion (max 60 seconds)
            gamma_url = None
            for _ in range(30):
                await asyncio.sleep(2)

                status_response = requests.get(
                    f"https://public-api.gamma.app/v1.0/generations/{generation_id}",
                    headers=headers,
                    timeout=15
                )

                if status_response.status_code == 200:
                    status_data = status_response.json()
                    status = status_data.get('status', '')

                    if status == 'completed' or status_data.get('gammaUrl'):
                        gamma_url = status_data.get('gammaUrl') or status_data.get('url')
                        break
                    elif status == 'failed':
                        return await self.generate_ppt(content)  # Fallback

            if gamma_url:
                return {
                    "success": True,
                    "presentation_url": gamma_url,
                    "presentation_id": generation_id,
                    "type": "gamma",
                    "message": "Professional presentation created with Gamma AI"
                }
            else:
                # Fallback to python-pptx
                return await self.generate_ppt(content)

        except Exception as e:
            print(f"Gamma API error: {str(e)}")
            # Fallback to python-pptx
            return await self.generate_ppt(content)

    async def _generate_html_presentation(self, content: Dict) -> Dict:
        """Fallback to HTML if PPT generation fails"""
        try:
            filename = f"presentation_{datetime.now().strftime('%Y%m%d_%H%M%S')}.html"
            filepath = os.path.join(self.deliverables_path, filename)

            title = content.get('title', 'Consulting Report')
            text = content.get('text', 'No content provided')

            slides = self._parse_content_to_slides(title, text)
            html_content = self._generate_presentation_html(title, slides)

            with open(filepath, 'w') as f:
                f.write(html_content)

            return {
                "success": True,
                "presentation_url": f"/api/deliverables/{filename}",
                "presentation_id": filename,
                "type": "html"
            }
        except Exception as e:
            return {"success": False, "error": str(e)}

    def _parse_content_to_slides(self, title: str, text: str) -> List[Dict]:
        """Parse text content into slide format"""
        slides = [{"title": title, "content": ["Executive Summary"], "type": "title"}]

        lines = text.split('\n')
        current_section = None
        current_content = []

        for line in lines:
            line = line.strip()
            if not line:
                continue

            if line.startswith('Problem:') or line.startswith('Challenge:'):
                if current_section:
                    slides.append({"title": current_section, "content": current_content, "type": "content"})
                current_section = "The Challenge"
                current_content = [line.replace('Problem:', '').replace('Challenge:', '').strip()]
            elif line.startswith('Consensus:') or line.startswith('Summary:'):
                if current_section:
                    slides.append({"title": current_section, "content": current_content, "type": "content"})
                current_section = "Executive Summary"
                current_content = [line.replace('Consensus:', '').replace('Summary:', '').strip()]
            elif line.startswith('Recommendation') or line.startswith('Next Step'):
                if current_section:
                    slides.append({"title": current_section, "content": current_content, "type": "content"})
                current_section = "Recommendations"
                current_content = []
            elif line.startswith(('-', '•', '*')) or (len(line) > 2 and line[0].isdigit() and line[1] in '.):'):
                clean = line.lstrip('-•*0123456789.)').strip()
                if clean:
                    current_content.append(clean)
            elif current_section:
                current_content.append(line)

        if current_section and current_content:
            slides.append({"title": current_section, "content": current_content, "type": "content"})

        slides.append({"title": "Next Steps", "content": ["Review recommendations", "Schedule follow-up", "Begin implementation"], "type": "content"})

        return slides

    def _generate_presentation_html(self, title: str, slides: List[Dict]) -> str:
        """Generate professional HTML presentation"""
        slides_html = ""
        for i, slide in enumerate(slides):
            slide_type = slide.get('type', 'content')

            if slide_type == 'title':
                slides_html += f'''
                <div class="slide title-slide" data-slide="{i}">
                    <h1>{slide['title']}</h1>
                    <p class="subtitle">{slide['content'][0] if slide['content'] else ''}</p>
                </div>
                '''
            else:
                content_html = ""
                for item in slide.get('content', []):
                    content_html += f"<li>{item}</li>\n"

                slides_html += f'''
                <div class="slide" data-slide="{i}">
                    <h2>{slide['title']}</h2>
                    <ul>{content_html}</ul>
                </div>
                '''

        return f'''<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <title>{title}</title>
    <style>
        * {{ margin: 0; padding: 0; box-sizing: border-box; }}
        body {{ font-family: 'Inter', sans-serif; background: linear-gradient(135deg, #1e293b 0%, #0f172a 100%); color: white; min-height: 100vh; }}
        .presentation {{ max-width: 1200px; margin: 0 auto; padding: 40px 20px; }}
        .slide {{ background: rgba(255,255,255,0.05); border-radius: 24px; padding: 60px; margin-bottom: 40px; min-height: 500px; }}
        .title-slide {{ text-align: center; background: linear-gradient(135deg, rgba(59,130,246,0.2), rgba(147,51,234,0.2)); }}
        .title-slide h1 {{ font-size: 3.5rem; background: linear-gradient(135deg, #60a5fa, #a78bfa); -webkit-background-clip: text; -webkit-text-fill-color: transparent; }}
        h2 {{ font-size: 2.5rem; color: #60a5fa; margin-bottom: 40px; }}
        li {{ font-size: 1.25rem; padding: 15px 0 15px 40px; position: relative; }}
        li::before {{ content: '→'; position: absolute; left: 0; color: #60a5fa; }}
    </style>
</head>
<body>
    <div class="presentation">{slides_html}</div>
</body>
</html>'''
